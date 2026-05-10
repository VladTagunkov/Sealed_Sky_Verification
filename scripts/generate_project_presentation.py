#!/usr/bin/env python3
"""Generate project overview PowerPoint (run: python3 scripts/generate_project_presentation.py)."""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "Trusted_Edge_Oracle_Sealed_Sky.pptx"


def add_title_slide(prs, title: str, subtitle: str) -> None:
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def add_bullets(prs, title: str, lines: list[str]) -> None:
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    for i, line in enumerate(lines):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(20)


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "Trusted Edge Oracle + Sealed Sky",
        "Hardware-trusted inference signing · timelock encryption · ENS / NameStone\nETH Prague 2026 · Space Computers",
    )

    add_bullets(
        prs,
        "What this project is",
        [
            "End-to-end story: prove inference came from secure hardware, verify off- and on-chain, optionally publish validation history to ENS.",
            "Monorepo with two user-facing apps that share branding and deep links.",
            "Built for hackathon demos: USB Armory MK II (ARM TrustZone) + Node backend + static verification UI + Vite/React Sealed Sky.",
        ],
    )

    add_bullets(
        prs,
        "Two parts — two dev servers",
        [
            "Part 1 — Trusted Edge Oracle: Node (Express) on localhost:3000 — API, HMAC verification, Sepolia relayer, static UI at /ui/.",
            "Part 2 — Sealed Sky: Vite + React on localhost:5173 — drand timelock + SpaceComputer cTRNG; optional NameStone capsules.",
            "Run both terminals when demoing the full flow; navigation links tie the apps together.",
        ],
    )

    add_bullets(
        prs,
        "Trusted Edge — hardware & applet",
        [
            "USB Armory MK II runs a Go/TamaGo Trusted OS + Rust #![no_std] applet in Secure World.",
            "Host talks over USB CDC-ECM: bridge at 10.0.0.1:4000 (newline-delimited JSON).",
            "Applet implements SignInference: HMAC-SHA256 over the full input string; hex output as proof.",
            "Hot-swap: make applet + upload ELF — no full SD reflash for Rust-only changes.",
        ],
    )

    add_bullets(
        prs,
        "Trusted Edge — backend & verification",
        [
            "Express server verifies proofs (canonical + legacy HMAC formats for older builds).",
            "Invalid proof returns HTTP 200 with success: false — clear UX for the verification UI.",
            "Accepted/rejected attempts logged to JSON history; optional NameStone publish (full history vs last event).",
            "Read path prefers NameStone HTTP API for robust ENS text resolution.",
        ],
    )

    add_bullets(
        prs,
        "On-chain — TrustedEdgeOracle (Foundry)",
        [
            "Solidity oracle on Sepolia: relayer submits attestations after backend verification.",
            "submitInference(deviceId, result, timestamp, imageHash, payloadHash) — only authorized relayers.",
            "Deploy with forge script; wire contract address into backend server.js.",
        ],
    )

    add_bullets(
        prs,
        "Sealed Sky — timelock & ENS",
        [
            "drand quicknet: non-interactive timelock (BLS12-381 IBE) in the browser.",
            "cTRNG mode: commit–reveal bound to future cosmic randomness witness from IPFS.",
            "ENS Phase 1: wallet connect, primary name, to/from metadata — no gas.",
            "ENS Phase 2: NameStone subdomains per capsule — envelope + unlock metadata on-chain readable text records.",
        ],
    )

    add_bullets(
        prs,
        "Architecture snapshot",
        [
            "Armory (10.0.0.1) ←→ host scripts (armory-link, nc, bun upload) ←→ backend :3000 ←→ Sepolia oracle.",
            "Sealed Sky static/Vite app — no server required for crypto; NameStone optional for publish.",
            "Firmware rebuild: ./docker/build.sh → flash SD; applet iteration: make applet + upload.",
        ],
    )

    add_bullets(
        prs,
        "Demo checklist",
        [
            "Bring up Armory USB; ./scripts/armory-link.sh; probe bridge with JSON-RPC.",
            "Open localhost:3000/ui/ for Seal Sky verification; localhost:5173 for Sealed Sky.",
            "Configure RPC + relayer key + optional NameStone in env files (never commit secrets).",
        ],
    )

    add_bullets(
        prs,
        "References",
        [
            "GoTEE — github.com/usbarmory/GoTEE",
            "USB Armory wiki — github.com/usbarmory/usbarmory/wiki",
            "ENS docs — docs.ens.domains",
            "NameStone — namestone.com",
        ],
    )

    prs.save(OUT)
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    main()
